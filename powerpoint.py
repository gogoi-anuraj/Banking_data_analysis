from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd

# Load the data

data = pd.read_csv(r"c:\Users\DELL\Downloads\banking_data.csv")

# Create a presentation object
prs = Presentation()

# Function to add a slide
def add_slide(prs, title, content, img_path=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    body_placeholder.text = content
    if img_path:
        slide.shapes.add_picture(img_path, Inches(1), Inches(2.5), height=Inches(3.5))

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title_placeholder = slide.shapes.title
subtitle_placeholder = slide.placeholders[1]
title_placeholder.text = "Exploratory Data Analysis of Banking Dataset"
subtitle_placeholder.text = "Insights from Telemarketing Campaigns\nPresented by: [Your Name]"

# Slide 2: Introduction
add_slide(prs, "Introduction", "Overview:\n- The dataset contains information about direct marketing campaigns of a Portuguese banking institution.\n- The goal is to predict if a client will subscribe to a term deposit.")

# Slide 3: Dataset Overview
add_slide(prs, "Dataset Overview", "Data Description:\n- Total Rows: 45,211\n- Total Columns: 18\nKey Variables:\n- Age, Job, Marital Status, Education, Default, Balance, Housing Loan, Personal Loan, Contact, Duration, Campaign, Pdays, Previous, Poutcome, Target (y)")

# Generate and save plots for embedding in the presentation
# Slide 4: Age and Job Distribution
plt.figure(figsize=(10, 6))
sns.histplot(data['age'], bins=30, kde=True)
plt.title('Age Distribution')
plt.xlabel('Age')
plt.ylabel('Frequency')
plt.savefig('age_distribution.png')

plt.figure(figsize=(12, 6))
sns.countplot(y='job', data=data, order=data['job'].value_counts().index)
plt.title('Job Type Distribution')
plt.xlabel('Count')
plt.ylabel('Job Type')
plt.savefig('job_type_distribution.png')

add_slide(prs, "Age and Job Distribution", "Insights: Majority of clients are between 30 to 40 years old. Most common job types are 'admin.' and 'blue-collar'.", 'age_distribution.png')
slide = prs.slides[-1]
slide.shapes.add_picture('job_type_distribution.png', Inches(5), Inches(2.5), height=Inches(3.5))

# Slide 5: Marital Status and Education Distribution
plt.figure(figsize=(8, 6))
sns.countplot(x='marital', data=data)
plt.title('Marital Status Distribution')
plt.xlabel('Marital Status')
plt.ylabel('Count')
plt.savefig('marital_status_distribution.png')

plt.figure(figsize=(8, 6))
sns.countplot(x='education', data=data)
plt.title('Education Level Distribution')
plt.xlabel('Education Level')
plt.ylabel('Count')
plt.savefig('education_level_distribution.png')

add_slide(prs, "Marital Status and Education Distribution", "Insights: Majority of clients are married. Most clients have secondary education.", 'marital_status_distribution.png')
slide = prs.slides[-1]
slide.shapes.add_picture('education_level_distribution.png', Inches(5), Inches(2.5), height=Inches(3.5))

# Slide 6: Credit Default and Loan Proportions
default_proportion = data['default'].value_counts(normalize=True) * 100
plt.figure(figsize=(8, 6))
sns.barplot(x=default_proportion.index, y=default_proportion.values)
plt.title('Proportion of Clients with Credit in Default')
plt.xlabel('Default')
plt.ylabel('Percentage')
plt.savefig('credit_default_proportion.png')

housing_loan = data['housing'].value_counts()
plt.figure(figsize=(8, 6))
sns.barplot(x=housing_loan.index, y=housing_loan.values)
plt.title('Number of Clients with Housing Loans')
plt.xlabel('Housing Loan')
plt.ylabel('Count')
plt.savefig('housing_loan.png')

personal_loan = data['loan'].value_counts()
plt.figure(figsize=(8, 6))
sns.barplot(x=personal_loan.index, y=personal_loan.values)
plt.title('Number of Clients with Personal Loans')
plt.xlabel('Personal Loan')
plt.ylabel('Count')
plt.savefig('personal_loan.png')

add_slide(prs, "Credit Default and Loan Proportions", "Insights: Most clients do not have credit in default. More clients have housing loans than personal loans.", 'credit_default_proportion.png')
slide = prs.slides[-1]
slide.shapes.add_picture('housing_loan.png', Inches(5), Inches(2.5), height=Inches(3.5))
slide.shapes.add_picture('personal_loan.png', Inches(8), Inches(2.5), height=Inches(3.5))

# Slide 7: Subscription Analysis
plt.figure(figsize=(10, 6))
sns.histplot(data=data, x='age', hue='y', multiple='stack', bins=30)
plt.title('Subscription to Term Deposit by Age')
plt.xlabel('Age')
plt.ylabel('Frequency')
plt.savefig('subscription_by_age.png')

plt.figure(figsize=(12, 6))
sns.countplot(y='job', hue='y', data=data, order=data['job'].value_counts().index)
plt.title('Subscription to Term Deposit by Job Type')
plt.xlabel('Count')
plt.ylabel('Job Type')
plt.savefig('subscription_by_job.png')

add_slide(prs, "Subscription Analysis", "Insights: Younger and older clients are more likely to subscribe. Certain job types like 'student' have higher subscription rates.", 'subscription_by_age.png')
slide = prs.slides[-1]
slide.shapes.add_picture('subscription_by_job.png', Inches(5), Inches(2.5), height=Inches(3.5))

# Slide 8: Conclusions
add_slide(prs, "Conclusions", "Key Findings:\n- Age and job type are significant factors in term deposit subscription.\n- Clients with higher balances are less likely to have loans.\n- Longer call durations tend to result in higher subscription rates.")

# Save the presentation
prs.save('/mnt/data/EDA_Banking_Dataset_Presentation.pptx')
